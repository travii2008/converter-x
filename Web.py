"""
Universal Converter - Flet Edition (Desktop + Android + Web)
Place in your Website folder alongside config.py and utils.py

Install:  pip install flet
Run:      python Web.py
Web:      flet build web --module-name Web
Android:  flet build apk --module-name Web
"""

import flet as ft
from pathlib import Path
import sys, threading, time

sys.path.insert(0, str(Path(__file__).parent))
import config, utils

# ── Colours ───────────────────────────────────────────────────────────────────
BG      = "#0B0D15"; SURFACE = "#10131C"; CARD = "#161B27"; CARD2 = "#1C2133"
BORDER  = "#252D3F"; ACCENT  = "#6C63FF"; ACCENT2 = "#4F8EF7"
GREEN   = "#22C55E"; RED     = "#EF4444"; ORANGE = "#F97316"; YELLOW = "#EAB308"
TEXT    = "#E8ECF5"; TEXT2   = "#8892A4"; TEXT3  = "#3D4A60"

# ── Shared helpers ────────────────────────────────────────────────────────────
def txt(s, size=12, color=TEXT, bold=False, mono=False, expand=False):
    return ft.Text(s, size=size, color=color,
                   weight=ft.FontWeight.BOLD if bold else None,
                   font_family="Consolas" if mono else None,
                   expand=expand)

def card(content, color=CARD, padding=16, radius=12, margin_bottom=12):
    return ft.Container(content=content, bgcolor=color, border_radius=radius,
                        padding=padding, margin=ft.margin.only(bottom=margin_bottom))

def label(s):
    return ft.Text(s, size=10, color=TEXT3)

def accent_btn(text, on_click, expand=True, icon=None):
    return ft.ElevatedButton(
        text=text, on_click=on_click, icon=icon,
        bgcolor=ACCENT, color="white",
        style=ft.ButtonStyle(shape=ft.RoundedRectangleBorder(radius=8)),
        expand=expand)

def secondary_btn(text, on_click, expand=False, icon=None):
    return ft.OutlinedButton(
        text=text, on_click=on_click, icon=icon,
        style=ft.ButtonStyle(
            color=TEXT2,
            side=ft.BorderSide(1, BORDER),
            shape=ft.RoundedRectangleBorder(radius=8)),
        expand=expand)

def browse_box(label_ref, on_click):
    """Reusable file browse box."""
    return ft.Container(
        content=ft.Column([
            ft.Text("📂", size=28),
            label_ref,
        ], horizontal_alignment=ft.CrossAxisAlignment.CENTER,
           spacing=6),
        bgcolor=CARD2, border_radius=10, padding=20,
        alignment=ft.alignment.center,
        on_click=on_click, ink=True)

def mk_progress():
    """Returns (progress_bar, pct_label, status_text) tuple."""
    bar  = ft.ProgressBar(value=0, bgcolor=BORDER, color=ACCENT, height=5, expand=True)
    pct  = ft.Text("", size=10, color=ACCENT, width=38, text_align=ft.TextAlign.RIGHT)
    stat = ft.Text("", size=11, color=TEXT2)
    return bar, pct, stat

def progress_row(bar, pct):
    return ft.Row([bar, pct], vertical_alignment=ft.CrossAxisAlignment.CENTER, spacing=6)


# ══════════════════════════════════════════════════════════════════════════════
#  UNITS PAGE
# ══════════════════════════════════════════════════════════════════════════════
def build_units_page(page: ft.Page):
    cats = list(config.UNIT_CATS.keys())

    def get_units(cat):
        d = config.UNIT_CATS[cat]
        return d['units'] if '__temp__' in d else list(d.keys())

    result_lbl = ft.Text("—", size=20, color=ACCENT, weight=ft.FontWeight.BOLD)
    qref_row   = ft.Row(wrap=True, spacing=6)
    val_field  = ft.TextField(hint_text="Enter value", bgcolor=CARD2, border_color=BORDER,
                               color=TEXT, text_size=14, expand=True,
                               keyboard_type=ft.KeyboardType.NUMBER)
    cat_dd  = ft.Dropdown(bgcolor=CARD2, border_color=BORDER, color=TEXT, text_size=12,
                           options=[ft.dropdown.Option(c) for c in cats], value=cats[0])
    from_dd = ft.Dropdown(bgcolor=CARD2, border_color=BORDER, color=TEXT, text_size=12, expand=True)
    to_dd   = ft.Dropdown(bgcolor=CARD2, border_color=BORDER, color=TEXT, text_size=12, expand=True)

    def calc(e=None):
        raw = (val_field.value or "").strip()
        if not raw:
            result_lbl.value = "—"; result_lbl.color = ACCENT
            qref_row.controls.clear(); page.update(); return
        try:
            v = float(raw)
        except ValueError:
            result_lbl.value = "Invalid"; result_lbl.color = RED
            page.update(); return

        d   = config.UNIT_CATS[cat_dd.value]
        res = utils.cvt_unit(v, from_dd.value, to_dd.value, d)
        result_lbl.value = utils.fmt_num(res) if res is not None else "Error"
        result_lbl.color = ACCENT

        units = get_units(cat_dd.value)
        qref_row.controls.clear()
        cnt = 0
        for u in units:
            if cnt >= 6 or u == from_dd.value: continue
            r = utils.cvt_unit(v, from_dd.value, u, d)
            if r is not None:
                qref_row.controls.append(
                    ft.Container(
                        content=ft.Text(f"{utils.fmt_num(r)} {u}", size=10, color=TEXT2),
                        bgcolor=CARD2, border_radius=6,
                        padding=ft.padding.symmetric(4, 8)))
                cnt += 1
        page.update()

    def on_cat(e=None):
        units = get_units(cat_dd.value)
        from_dd.options = [ft.dropdown.Option(u) for u in units]
        to_dd.options   = [ft.dropdown.Option(u) for u in units]
        from_dd.value   = units[0]
        to_dd.value     = units[1] if len(units) > 1 else units[0]
        calc()

    def swap(e):
        from_dd.value, to_dd.value = to_dd.value, from_dd.value
        try: val_field.value = result_lbl.value.replace(",", "")
        except Exception: pass
        calc()

    cat_dd.on_change  = on_cat
    from_dd.on_change = calc
    to_dd.on_change   = calc
    val_field.on_change = calc
    on_cat()

    return ft.Column([
        txt("Unit Converter", size=15, bold=True),
        txt("Convert between all common units", size=11, color=TEXT2),
        ft.Container(height=10),
        card(ft.Column([label("Category"), ft.Container(height=4), cat_dd])),
        card(ft.Column([
            ft.Row([
                ft.Column([label("From"), ft.Container(height=4), val_field,
                           ft.Container(height=4), from_dd],
                          expand=True, spacing=0),
                ft.Column([
                    ft.Container(height=24),
                    ft.IconButton(ft.Icons.SWAP_HORIZ, icon_color=ACCENT, on_click=swap,
                                  bgcolor=CARD2, icon_size=22),
                ]),
                ft.Column([
                    label("To"), ft.Container(height=4),
                    ft.Container(content=result_lbl, bgcolor=CARD2, border_radius=8,
                                 padding=ft.padding.symmetric(10, 12), expand=True),
                    ft.Container(height=4), to_dd,
                ], expand=True, spacing=0),
            ], alignment=ft.MainAxisAlignment.SPACE_BETWEEN),
        ])),
        card(ft.Column([label("Quick Reference"), ft.Container(height=6), qref_row]),
             color=CARD2),
    ], scroll=ft.ScrollMode.AUTO, expand=True, spacing=0)


# ══════════════════════════════════════════════════════════════════════════════
#  CURRENCY PAGE
# ══════════════════════════════════════════════════════════════════════════════
def build_currency_page(page: ft.Page):
    MAJORS     = ["USD","EUR","GBP","JPY","CAD","AUD","CHF","CNY","INR","MXN",
                  "BRL","KRW","SGD","AED","TRY","NOK","SEK","PLN"]
    currencies = sorted(config.CURRENCY_RATES.keys())

    result_lbl = ft.Text("—", size=20, color=ACCENT, weight=ft.FontWeight.BOLD)
    rate_lbl   = ft.Text("", size=10, color=TEXT3)
    comp_grid  = ft.GridView(runs_count=6, spacing=4, run_spacing=4,
                              expand=False, max_extent=130)
    amt_field  = ft.TextField(value="1", bgcolor=CARD2, border_color=BORDER,
                               color=TEXT, text_size=14, expand=True,
                               keyboard_type=ft.KeyboardType.NUMBER)
    from_dd = ft.Dropdown(bgcolor=CARD2, border_color=BORDER, color=TEXT, text_size=11,
                           expand=True, editable=True,
                           options=[ft.dropdown.Option(c) for c in currencies], value="USD")
    to_dd   = ft.Dropdown(bgcolor=CARD2, border_color=BORDER, color=TEXT, text_size=11,
                           expand=True, editable=True,
                           options=[ft.dropdown.Option(c) for c in currencies], value="EUR")

    def calc(e=None):
        # Auto-uppercase
        if from_dd.value: from_dd.value = from_dd.value.upper()
        if to_dd.value:   to_dd.value   = to_dd.value.upper()

        try:
            amt = float(amt_field.value or "0")
        except ValueError:
            result_lbl.value = "—"; page.update(); return

        fc, tc = from_dd.value, to_dd.value
        if fc not in config.CURRENCY_RATES or tc not in config.CURRENCY_RATES:
            page.update(); return

        res  = utils.cvt_currency(amt, fc, tc)
        rate = config.CURRENCY_RATES[tc] / config.CURRENCY_RATES[fc]
        result_lbl.value = f"{res:,.4f}".rstrip("0").rstrip(".")
        rate_lbl.value   = f"1 {fc} = {rate:.6f} {tc}"

        comp_grid.controls.clear()
        for cur in MAJORS:
            r   = utils.cvt_currency(amt, fc, cur)
            val = "—" if cur == fc or r is None else f"{r:,.2f}"
            comp_grid.controls.append(
                ft.Container(bgcolor=BORDER, border_radius=6,
                             padding=ft.padding.symmetric(5, 8),
                             content=ft.Row([
                                 ft.Text(cur,  size=10, color=TEXT2),
                                 ft.Text(val,  size=10,
                                         color=TEXT3 if cur == fc else ACCENT,
                                         weight=ft.FontWeight.BOLD),
                             ], alignment=ft.MainAxisAlignment.SPACE_BETWEEN)))
        page.update()

    def swap(e):
        from_dd.value, to_dd.value = to_dd.value, from_dd.value
        calc()

    amt_field.on_change = calc
    from_dd.on_change   = calc
    to_dd.on_change     = calc
    calc()

    return ft.Column([
        txt("Currency Converter", size=15, bold=True),
        txt(f"{len(currencies)} currencies · live rates", size=11, color=TEXT2),
        ft.Container(height=10),
        card(ft.Column([
            ft.Row([
                ft.Column([label("Amount"), ft.Container(height=4),
                           amt_field, ft.Container(height=4), from_dd],
                          expand=True, spacing=0),
                ft.Column([
                    ft.Container(height=24),
                    ft.IconButton(ft.Icons.SWAP_HORIZ, icon_color=ACCENT,
                                  on_click=swap, bgcolor=CARD2, icon_size=22),
                ]),
                ft.Column([
                    label("Result"), ft.Container(height=4),
                    ft.Container(content=result_lbl, bgcolor=CARD2, border_radius=8,
                                 padding=ft.padding.symmetric(10, 12), expand=True),
                    ft.Container(height=4), to_dd,
                ], expand=True, spacing=0),
            ]),
            ft.Container(height=6),
            rate_lbl,
        ])),
        card(ft.Column([
            label("Quick Compare — major currencies"),
            ft.Container(height=6),
            comp_grid,
        ]), color=CARD2),
    ], scroll=ft.ScrollMode.AUTO, expand=True, spacing=0)


# ══════════════════════════════════════════════════════════════════════════════
#  AUDIO PAGE
# ══════════════════════════════════════════════════════════════════════════════
def build_audio_page(page: ft.Page):
    file_lbl  = ft.Text("Click to browse audio file", size=11, color=TEXT2)
    bar, pct_lbl, stat = mk_progress()
    selected  = {"file": None}

    fmt_group = ft.RadioGroup(
        content=ft.Row([ft.Radio(value=f, label=f.upper(), fill_color=ACCENT)
                        for f in config.AUDIO_FMTS], wrap=True),
        value="mp3")
    br_dd = ft.Dropdown(value="192k", bgcolor=CARD2, border_color=BORDER,
                         color=TEXT, text_size=11, expand=True,
                         options=[ft.dropdown.Option(b)
                                  for b in ["64k","96k","128k","160k","192k","256k","320k"]])
    sr_dd = ft.Dropdown(value="44100", bgcolor=CARD2, border_color=BORDER,
                         color=TEXT, text_size=11, expand=True,
                         options=[ft.dropdown.Option(s)
                                  for s in ["22050","44100","48000","96000"]])
    ch_dd = ft.Dropdown(value="2", bgcolor=CARD2, border_color=BORDER,
                         color=TEXT, text_size=11, expand=True,
                         options=[ft.dropdown.Option("1","Mono"),
                                  ft.dropdown.Option("2","Stereo")])

    def pick_result(e: ft.FilePickerResultEvent):
        if e.files:
            selected["file"] = e.files[0].path
            file_lbl.value   = e.files[0].name
            file_lbl.color   = TEXT
        page.update()

    picker = ft.FilePicker(on_result=pick_result)
    page.overlay.append(picker)

    def save_result(e: ft.FilePickerResultEvent):
        if not e.path: return
        fmt = fmt_group.value
        out = e.path if e.path.endswith(f".{fmt}") else e.path + f".{fmt}"
        stat.value = "Converting…"; stat.color = YELLOW
        bar.value  = 0; pct_lbl.value = "0%"; page.update()

        def run():
            def on_prog(p):
                bar.value      = p / 100
                pct_lbl.value  = f"{int(p)}%"
                pct_lbl.color  = GREEN if p >= 100 else ACCENT
                page.update()
            ok, msg = utils.run_ffmpeg_progress(
                ["-i", selected["file"], "-ar", sr_dd.value,
                 "-ac", ch_dd.value, "-b:a", br_dd.value, out],
                on_prog)
            def done():
                if ok:
                    stat.value = "✓ Saved!"; stat.color = GREEN
                    bar.value  = 1; pct_lbl.value = "100%"
                    utils.save_history({"type":"Audio",
                                        "input": Path(selected["file"]).name,
                                        "output": Path(out).name,
                                        "note": f"{fmt.upper()} · {br_dd.value}"})
                else:
                    stat.value = f"✗ {msg}"; stat.color = RED
                    bar.value  = 0; pct_lbl.value = ""
                page.update()
            page.run_task(lambda: done()) if hasattr(page, 'run_task') else done()
        threading.Thread(target=run, daemon=True).start()

    save_picker = ft.FilePicker(on_result=save_result)
    page.overlay.append(save_picker)

    def convert(e):
        if not selected["file"]:
            stat.value = "Select a file first"; stat.color = RED; page.update(); return
        fmt  = fmt_group.value
        stem = Path(selected["file"]).stem
        save_picker.save_file(file_name=f"{stem}.{fmt}", allowed_extensions=[fmt])

    return ft.Column([
        txt("Audio Converter", size=15, bold=True),
        ft.Text("⚠ FFmpeg required for conversion", size=10, color=ORANGE)
        if not config.HAS_FFMPEG else ft.Container(height=0),
        ft.Container(height=8),
        card(ft.Column([
            browse_box(file_lbl,
                       lambda _: picker.pick_files(
                           allowed_extensions=config.AUDIO_FMTS + config.VIDEO_FMTS)),
            ft.Container(height=8),
            label("Output Format"),
            fmt_group,
            ft.Container(height=6),
            ft.Row([
                ft.Column([label("Bitrate"),      ft.Container(height=4), br_dd], expand=True),
                ft.Column([label("Sample Rate"),  ft.Container(height=4), sr_dd], expand=True),
                ft.Column([label("Channels"),     ft.Container(height=4), ch_dd], expand=True),
            ], spacing=10),
            ft.Container(height=8),
            progress_row(bar, pct_lbl),
            ft.Container(height=4),
            stat,
            ft.Container(height=6),
            accent_btn("🎵 Convert Audio", convert),
        ])),
    ], scroll=ft.ScrollMode.AUTO, expand=True, spacing=0)


# ══════════════════════════════════════════════════════════════════════════════
#  VIDEO PAGE
# ══════════════════════════════════════════════════════════════════════════════
def build_video_page(page: ft.Page):
    file_lbl  = ft.Text("Click to browse video file", size=11, color=TEXT2)
    bar, pct_lbl, stat = mk_progress()
    selected  = {"file": None}

    fmt_group = ft.RadioGroup(
        content=ft.Row([ft.Radio(value=f, label=f.upper(), fill_color=ACCENT)
                        for f in ["mp4","avi","mov","mkv","webm","flv","wmv"]], wrap=True),
        value="mp4")
    crf_dd = ft.Dropdown(value="23", bgcolor=CARD2, border_color=BORDER,
                          color=TEXT, text_size=11, expand=True,
                          options=[ft.dropdown.Option("18","18 (Best)"),
                                   ft.dropdown.Option("20","20"),
                                   ft.dropdown.Option("23","23 (Default)"),
                                   ft.dropdown.Option("26","26"),
                                   ft.dropdown.Option("30","30 (Small)")])
    res_dd = ft.Dropdown(value="", bgcolor=CARD2, border_color=BORDER,
                          color=TEXT, text_size=11, expand=True,
                          options=[ft.dropdown.Option("","Original"),
                                   ft.dropdown.Option("1920x1080","1080p"),
                                   ft.dropdown.Option("1280x720","720p"),
                                   ft.dropdown.Option("854x480","480p"),
                                   ft.dropdown.Option("640x360","360p")])
    fps_dd = ft.Dropdown(value="", bgcolor=CARD2, border_color=BORDER,
                          color=TEXT, text_size=11, expand=True,
                          options=[ft.dropdown.Option("","Original"),
                                   ft.dropdown.Option("60","60"),
                                   ft.dropdown.Option("30","30"),
                                   ft.dropdown.Option("24","24")])

    def pick_result(e: ft.FilePickerResultEvent):
        if e.files:
            selected["file"] = e.files[0].path
            file_lbl.value   = e.files[0].name
            file_lbl.color   = TEXT
        page.update()

    picker = ft.FilePicker(on_result=pick_result)
    page.overlay.append(picker)

    def _run(args, out, note):
        stat.value = "Converting…"; stat.color = YELLOW
        bar.value  = 0; pct_lbl.value = "0%"; page.update()
        def run():
            def on_prog(p):
                bar.value     = p / 100
                pct_lbl.value = f"{int(p)}%"
                pct_lbl.color = GREEN if p >= 100 else ACCENT
                page.update()
            ok, msg = utils.run_ffmpeg_progress(args, on_prog)
            def done():
                if ok:
                    stat.value = "✓ Saved!"; stat.color = GREEN
                    bar.value  = 1; pct_lbl.value = "100%"
                    utils.save_history({"type":"Video",
                                        "input": Path(selected["file"]).name,
                                        "output": Path(out).name,
                                        "note": note})
                else:
                    stat.value = f"✗ {msg}"; stat.color = RED
                    bar.value  = 0; pct_lbl.value = ""
                page.update()
            done()
        threading.Thread(target=run, daemon=True).start()

    def save_video(e: ft.FilePickerResultEvent):
        if not e.path: return
        fmt  = fmt_group.value
        out  = e.path if e.path.endswith(f".{fmt}") else e.path + f".{fmt}"
        args = ["-i", selected["file"], "-crf", crf_dd.value]
        if res_dd.value: args += ["-vf", f"scale={res_dd.value}"]
        if fps_dd.value: args += ["-r", fps_dd.value]
        args.append(out)
        _run(args, out, f"{fmt.upper()} · CRF {crf_dd.value}")

    def save_audio(e: ft.FilePickerResultEvent):
        if not e.path: return
        out = e.path if e.path.endswith(".mp3") else e.path + ".mp3"
        _run(["-i", selected["file"], "-vn", "-ab", "192k", out], out, "MP3 · 192k")

    save_v = ft.FilePicker(on_result=save_video)
    save_a = ft.FilePicker(on_result=save_audio)
    page.overlay.extend([save_v, save_a])

    def do_convert(e):
        if not selected["file"]:
            stat.value = "Select a file first"; stat.color = RED; page.update(); return
        fmt = fmt_group.value
        save_v.save_file(file_name=Path(selected["file"]).stem + f".{fmt}",
                          allowed_extensions=[fmt])

    def do_extract(e):
        if not selected["file"]:
            stat.value = "Select a file first"; stat.color = RED; page.update(); return
        save_a.save_file(file_name=Path(selected["file"]).stem + ".mp3",
                          allowed_extensions=["mp3"])

    return ft.Column([
        txt("Video Converter", size=15, bold=True),
        ft.Text("⚠ FFmpeg required for conversion", size=10, color=ORANGE)
        if not config.HAS_FFMPEG else ft.Container(height=0),
        ft.Container(height=8),
        card(ft.Column([
            browse_box(file_lbl,
                       lambda _: picker.pick_files(allowed_extensions=config.VIDEO_FMTS)),
            ft.Container(height=8),
            label("Output Format"),
            fmt_group,
            ft.Container(height=6),
            ft.Row([
                ft.Column([label("Quality (CRF)"), ft.Container(height=4), crf_dd], expand=True),
                ft.Column([label("Resolution"),    ft.Container(height=4), res_dd], expand=True),
                ft.Column([label("FPS"),           ft.Container(height=4), fps_dd], expand=True),
            ], spacing=10),
            ft.Container(height=8),
            progress_row(bar, pct_lbl),
            ft.Container(height=4),
            stat,
            ft.Container(height=6),
            ft.Row([
                accent_btn("🎬 Convert Video",     do_convert),
                secondary_btn("🎵 Extract Audio",  do_extract, expand=True),
            ], spacing=10),
        ])),
    ], scroll=ft.ScrollMode.AUTO, expand=True, spacing=0)


# ══════════════════════════════════════════════════════════════════════════════
#  IMAGE PAGE
# ══════════════════════════════════════════════════════════════════════════════
def build_image_page(page: ft.Page):
    file_lbl  = ft.Text("Click to browse image(s)", size=11, color=TEXT2)
    bar, pct_lbl, stat = mk_progress()
    selected  = {"files": []}
    file_list = ft.Column(spacing=4)

    fmt_group = ft.RadioGroup(
        content=ft.Row([ft.Radio(value=f, label=f.upper(), fill_color=ACCENT)
                        for f in config.IMG_OUT], wrap=True),
        value="png")
    scale_field = ft.TextField(value="100", bgcolor=CARD2, border_color=BORDER,
                                color=TEXT, text_size=12,
                                keyboard_type=ft.KeyboardType.NUMBER,
                                width=90, suffix_text="%")

    def render_list():
        file_list.controls.clear()
        for i, fp in enumerate(selected["files"]):
            p = Path(fp)
            try:   sz = f"{p.stat().st_size // 1024}KB"
            except: sz = ""
            file_list.controls.append(
                ft.Container(bgcolor=CARD2, border_radius=6, padding=8,
                             content=ft.Row([
                                 ft.Text("🖼", size=14),
                                 ft.Text(p.name, size=10, color=TEXT, expand=True),
                                 ft.Text(sz, size=10, color=TEXT3),
                                 ft.IconButton(ft.Icons.CLOSE, icon_size=14,
                                               icon_color=TEXT3,
                                               on_click=lambda _, idx=i: rm(idx)),
                             ])))
        page.update()

    def rm(i):
        selected["files"].pop(i)
        if not selected["files"]:
            file_lbl.value = "Click to browse image(s)"
            file_lbl.color = TEXT2
        render_list()

    def pick_result(e: ft.FilePickerResultEvent):
        if e.files:
            selected["files"] = [f.path for f in e.files]
            file_lbl.value = (e.files[0].name if len(e.files) == 1
                              else f"{len(e.files)} files selected")
            file_lbl.color = TEXT
            render_list()
        page.update()

    def save_result(e: ft.FilePickerResultEvent):
        if not e.path: return
        out_dir = e.path
        total   = len(selected["files"])
        stat.value = f"Converting {total} image(s)…"
        stat.color = YELLOW; bar.value = 0; pct_lbl.value = "0%"; page.update()

        def run():
            from PIL import Image
            fmt   = fmt_group.value
            scale = int(scale_field.value or "100") / 100
            done_c = 0
            for i, fp in enumerate(selected["files"]):
                try:
                    img = Image.open(fp)
                    if scale != 1.0:
                        img = img.resize((max(1, int(img.width*scale)),
                                          max(1, int(img.height*scale))), Image.LANCZOS)
                    if fmt in ("pdf","jpg","jpeg") and img.mode in ("RGBA","P","LA"):
                        img = img.convert("RGB")
                    dst = str(Path(out_dir) / (Path(fp).stem + f".{fmt}"))
                    img.save(dst)
                    done_c += 1
                except Exception:
                    pass
                pct = ((i + 1) / total) * 100
                bar.value     = pct / 100
                pct_lbl.value = f"{int(pct)}%"
                page.update()

            stat.value = f"✓ {done_c}/{total} converted"
            stat.color = GREEN; bar.value = 1; pct_lbl.value = "100%"
            utils.save_history({"type":"Image",
                                 "input": f"{total} file(s)",
                                 "output": fmt.upper(),
                                 "note": f"Scale {scale_field.value}%"})
            page.update()

        threading.Thread(target=run, daemon=True).start()

    picker      = ft.FilePicker(on_result=pick_result)
    save_picker = ft.FilePicker(on_result=save_result)
    page.overlay.extend([picker, save_picker])

    def convert(e):
        if not selected["files"]:
            stat.value = "Select file(s) first"; stat.color = RED; page.update(); return
        save_picker.get_directory_path(dialog_title="Select output folder")

    return ft.Column([
        txt("Image Converter", size=15, bold=True),
        ft.Text("⚠ Pillow not installed", size=10, color=ORANGE)
        if not config.HAS_PIL else ft.Container(height=0),
        ft.Container(height=8),
        card(ft.Column([
            browse_box(file_lbl,
                       lambda _: picker.pick_files(
                           allow_multiple=True, allowed_extensions=config.IMG_IN)),
            file_list,
            ft.Container(height=8),
            label("Convert To"),
            fmt_group,
            ft.Container(height=6),
            ft.Row([
                label("Resize Scale (%)"),
                ft.Container(width=10),
                scale_field,
            ], vertical_alignment=ft.CrossAxisAlignment.CENTER),
            ft.Container(height=8),
            progress_row(bar, pct_lbl),
            ft.Container(height=4),
            stat,
            ft.Container(height=6),
            accent_btn("🖼 Convert Image(s)", convert),
        ])),
    ], scroll=ft.ScrollMode.AUTO, expand=True, spacing=0)


# ══════════════════════════════════════════════════════════════════════════════
#  PDF PAGE
# ══════════════════════════════════════════════════════════════════════════════
def build_pdf_page(page: ft.Page):
    stat = ft.Text("", size=11, color=TEXT2)
    bar, pct_lbl, _ = mk_progress()

    def log(msg, color=TEXT2):
        stat.value = msg; stat.color = color; page.update()

    def _run_thread(task_fn, success_msg, hist_entry=None):
        log("Working…", YELLOW); bar.value = 0.2; pct_lbl.value = ""; page.update()
        def run():
            ok, msg = task_fn()
            if ok:
                log(f"✓ {success_msg}", GREEN)
                bar.value = 1; pct_lbl.value = "100%"
                if hist_entry: utils.save_history(hist_entry)
            else:
                log(f"✗ {msg}", RED)
                bar.value = 0; pct_lbl.value = ""
            page.update()
        threading.Thread(target=run, daemon=True).start()

    # ── Pickers ───────────────────────────────────────────────────────────
    # Each operation gets its own picker pair (source + dest)
    src_store  = {"path": None}
    dest_store = {"path": None, "action": None}

    def mk_src_picker(exts, action_fn):
        """Returns a picker that stores the source then opens the save dialog."""
        def on_pick(e: ft.FilePickerResultEvent):
            if e.files:
                src_store["path"] = e.files[0].path
                action_fn(e.files[0].path)
        p = ft.FilePicker(on_result=on_pick)
        page.overlay.append(p)
        return lambda: p.pick_files(allowed_extensions=exts)

    def mk_dest_picker(default_ext, task_builder):
        """Returns (picker, open_fn). open_fn(src) opens the save dialog."""
        def on_save(e: ft.FilePickerResultEvent):
            if not e.path: return
            out = e.path if e.path.endswith(default_ext) else e.path + default_ext
            _run_thread(lambda: task_builder(src_store["path"], out),
                        f"Saved as {Path(out).name}",
                        {"type":"PDF","input":Path(src_store["path"]).name,
                         "output":Path(out).name,"note":f"→ {default_ext.upper()}"})
        p = ft.FilePicker(on_result=on_save)
        page.overlay.append(p)
        return p

    # ── To PDF ────────────────────────────────────────────────────────────
    def task_img_pdf(src, out):
        try:
            if config.HAS_PIL:
                from PIL import Image as _I
                img = _I.open(src)
                if img.mode in ("RGBA","P","LA"): img = img.convert("RGB")
                img.save(out, "PDF")
                return True, "OK"
            return False, "Pillow not installed"
        except Exception as e: return False, str(e)

    def task_office_pdf(src, out):
        ok, msg = utils.run_office_conversion(src, "pdf")
        if ok:
            gen = Path(src).with_suffix(".pdf")
            if gen.exists() and str(gen) != out:
                import shutil; shutil.move(str(gen), out)
        return ok, msg

    def task_text_pdf(src, out):
        try:
            from PIL import Image as _I, ImageDraw
            lines = Path(src).read_text(encoding="utf-8", errors="replace").split("\n")
            W, H, margin, lh = 794, 1123, 60, 20
            per_page = (H - 2*margin) // lh
            pages_list = []
            for i in range(0, max(1, len(lines)), per_page):
                img = _I.new("RGB", (W, H), "white")
                draw = ImageDraw.Draw(img)
                for j, line in enumerate(lines[i:i+per_page]):
                    draw.text((margin, margin + j*lh), line, fill="black")
                pages_list.append(img)
            pages_list[0].save(out, save_all=True, append_images=pages_list[1:])
            return True, "OK"
        except Exception as e: return False, str(e)

    # ── From PDF ──────────────────────────────────────────────────────────
    def task_pdf_word(src, out):
        try:
            if config.HAS_PDF2DOCX:
                from pdf2docx import Converter as CV
                cv = CV(src); cv.convert(out); cv.close()
                return True, "OK"
            elif config.HAS_PDFPLUMBER and config.HAS_DOCX:
                import pdfplumber
                from docx import Document
                doc = Document()
                with pdfplumber.open(src) as pdf:
                    for p in pdf.pages:
                        for line in (p.extract_text() or "").split("\n"):
                            doc.add_paragraph(line)
                        doc.add_page_break()
                doc.save(out)
                return True, "OK"
            return False, "Install pdf2docx:  pip install pdf2docx"
        except Exception as e: return False, str(e)

    def task_pdf_text(src, out):
        try:
            import pdfplumber
            with pdfplumber.open(src) as pdf:
                text = "\n\n--- Page Break ---\n\n".join(
                    p.extract_text() or "" for p in pdf.pages)
            Path(out).write_text(text, encoding="utf-8")
            return True, "OK"
        except Exception as e: return False, str(e)

    def task_pdf_excel(src, out):
        try:
            import pdfplumber, openpyxl
            wb = openpyxl.Workbook(); wb.remove(wb.active)
            with pdfplumber.open(src) as pdf:
                tc = 0
                for pn, p in enumerate(pdf.pages, 1):
                    for tn, table in enumerate(p.extract_tables(), 1):
                        tc += 1
                        ws = wb.create_sheet(f"P{pn}_T{tn}")
                        for row in table: ws.append(row)
            if tc == 0:
                ws = wb.create_sheet("Text")
                with pdfplumber.open(src) as pdf:
                    for p in pdf.pages:
                        for line in (p.extract_text() or "").split("\n"):
                            ws.append([line])
            wb.save(out); return True, "OK"
        except Exception as e: return False, str(e)

    def task_pdf_ppt(src, out):
        try:
            import pdfplumber
            from pptx import Presentation
            from pptx.util import Inches
            from PIL import Image as _I
            import tempfile, os
            prs    = Presentation()
            blank  = prs.slide_layouts[6]
            W, H   = prs.slide_width, prs.slide_height
            with pdfplumber.open(src) as pdf:
                for p in pdf.pages:
                    img = p.to_image(resolution=150)
                    tmp = tempfile.mktemp(suffix=".png")
                    img.save(tmp)
                    slide = prs.slides.add_slide(blank)
                    slide.shapes.add_picture(tmp, 0, 0, W, H)
                    os.remove(tmp)
            prs.save(out); return True, "OK"
        except Exception as e:
            return False, f"{e}  (pip install python-pptx pdfplumber)"

    # ── Wire pickers ──────────────────────────────────────────────────────
    p_img_pdf    = mk_dest_picker(".pdf", task_img_pdf)
    p_office_pdf = mk_dest_picker(".pdf", task_office_pdf)
    p_text_pdf   = mk_dest_picker(".pdf", task_text_pdf)
    p_pdf_word   = mk_dest_picker(".docx", task_pdf_word)
    p_pdf_text   = mk_dest_picker(".txt",  task_pdf_text)
    p_pdf_excel  = mk_dest_picker(".xlsx", task_pdf_excel)
    p_pdf_ppt    = mk_dest_picker(".pptx", task_pdf_ppt)

    def open_src_then_save(src_exts, dest_picker, dest_ext, stem_suffix=""):
        def on_src(e: ft.FilePickerResultEvent):
            if not e.files: return
            src_store["path"] = e.files[0].path
            stem = Path(e.files[0].path).stem
            dest_picker.save_file(file_name=f"{stem}{stem_suffix}{dest_ext}",
                                   allowed_extensions=[dest_ext.lstrip(".")])
        p = ft.FilePicker(on_result=on_src)
        page.overlay.append(p)
        return lambda: p.pick_files(allowed_extensions=src_exts)

    open_img_pdf    = open_src_then_save(config.IMG_IN,            p_img_pdf,    ".pdf")
    open_office_pdf = open_src_then_save(["docx","doc","xlsx","xls","pptx","ppt"], p_office_pdf, ".pdf")
    open_text_pdf   = open_src_then_save(["txt","md"],             p_text_pdf,   ".pdf")
    open_pdf_word   = open_src_then_save(["pdf"],                  p_pdf_word,   ".docx")
    open_pdf_text   = open_src_then_save(["pdf"],                  p_pdf_text,   ".txt")
    open_pdf_excel  = open_src_then_save(["pdf"],                  p_pdf_excel,  ".xlsx")
    open_pdf_ppt    = open_src_then_save(["pdf"],                  p_pdf_ppt,    ".pptx")

    # ── UI ────────────────────────────────────────────────────────────────
    def op_btn(label_text, on_click, color=None):
        return ft.ElevatedButton(
            text=label_text, on_click=lambda _: on_click(),
            bgcolor=color or BORDER, color=TEXT2,
            style=ft.ButtonStyle(shape=ft.RoundedRectangleBorder(radius=8)),
            expand=True)

    return ft.Column([
        txt("PDF Converter", size=15, bold=True),
        ft.Container(height=10),
        ft.Row([
            # LEFT — To PDF
            ft.Container(expand=True,
                content=card(ft.Column([
                    txt("⬆  Convert TO PDF", size=13, bold=True, color=ACCENT),
                    ft.Container(height=8),
                    op_btn("🖼  Image → PDF",   open_img_pdf),
                    ft.Container(height=4),
                    op_btn("📄  Word → PDF",    open_office_pdf),
                    ft.Container(height=4),
                    op_btn("📊  Excel → PDF",   open_office_pdf),
                    ft.Container(height=4),
                    op_btn("📑  PPT → PDF",     open_office_pdf),
                    ft.Container(height=4),
                    op_btn("📝  Text/MD → PDF", open_text_pdf),
                ]), margin_bottom=0)),
            ft.Container(width=12),
            # RIGHT — From PDF
            ft.Container(expand=True,
                content=card(ft.Column([
                    txt("⬇  Convert FROM PDF", size=13, bold=True, color=ACCENT2),
                    ft.Container(height=8),
                    op_btn("🖼  PDF → JPG",  open_pdf_text,  CARD2),
                    ft.Container(height=4),
                    op_btn("📄  PDF → Word", open_pdf_word,  CARD2),
                    ft.Container(height=4),
                    op_btn("📝  PDF → Text", open_pdf_text,  CARD2),
                    ft.Container(height=4),
                    op_btn("📊  PDF → Excel",open_pdf_excel, CARD2),
                    ft.Container(height=4),
                    op_btn("📑  PDF → PPT",  open_pdf_ppt,   CARD2),
                ]), margin_bottom=0)),
        ], vertical_alignment=ft.CrossAxisAlignment.START),
        ft.Container(height=10),
        progress_row(bar, pct_lbl),
        ft.Container(height=6),
        stat,
    ], scroll=ft.ScrollMode.AUTO, expand=True, spacing=0)


# ══════════════════════════════════════════════════════════════════════════════
#  ABOUT + HISTORY PAGE
# ══════════════════════════════════════════════════════════════════════════════
def build_about_page(page: ft.Page):
    deps = [
        ("FFmpeg",      "Audio/Video engine",       "ffmpeg.org",              config.HAS_FFMPEG),
        ("Pillow",      "Image conversion",          "pip install pillow",      config.HAS_PIL),
        ("pdfplumber",  "PDF reading",               "pip install pdfplumber",  config.HAS_PDFPLUMBER),
        ("python-docx", "Word documents",            "pip install python-docx", config.HAS_DOCX),
        ("pdf2docx",    "Better PDF → Word",         "pip install pdf2docx",    config.HAS_PDF2DOCX),
        ("openpyxl",    "Excel files",               "pip install openpyxl",    config.HAS_OPENPYXL),
        ("python-pptx", "PPT creation",              "pip install python-pptx", False),
        ("pywin32",     "MS Office (Windows only)",  "pip install pywin32",     config.HAS_OFFICE),
        ("LibreOffice", "Office → PDF fallback",     "libreoffice.org",         config.HAS_LIBREOFFICE),
    ]

    dep_rows = [
        ft.Row([
            ft.Text("✓" if ok else "✗", size=12,
                    color=GREEN if ok else RED, width=18),
            ft.Text(name, size=11, color=TEXT,
                    weight=ft.FontWeight.BOLD, width=120),
            ft.Text(desc, size=10, color=TEXT2, expand=True),
            ft.Text(cmd if not ok else "", size=10,
                    color=YELLOW, font_family="Consolas"),
        ])
        for name, desc, cmd, ok in deps
    ]

    # History
    history      = utils.load_history()
    hist_controls = []
    type_colors   = {"Audio":ACCENT,"Video":ACCENT2,"Image":GREEN,
                     "PDF":ORANGE,"Video → Audio":YELLOW}

    if not history:
        hist_controls.append(
            ft.Text("No conversions yet.", size=11, color=TEXT3))
    else:
        for entry in history:
            tc = type_colors.get(entry.get("type",""), TEXT3)
            ts = time.strftime("%b %d  %H:%M", time.localtime(entry.get("time",0)))
            hist_controls.append(
                ft.Container(bgcolor=BORDER, border_radius=6, padding=8,
                             margin=ft.margin.only(bottom=4),
                             content=ft.Row([
                                 ft.Text(entry.get("type","?"), size=10,
                                         color=tc, weight=ft.FontWeight.BOLD, width=90),
                                 ft.Text(f"{entry.get('input','?')} → {entry.get('output','?')}",
                                         size=10, color=TEXT, expand=True),
                                 ft.Column([
                                     ft.Text(entry.get("note",""), size=9, color=TEXT3),
                                     ft.Text(ts, size=9, color=TEXT3),
                                 ], horizontal_alignment=ft.CrossAxisAlignment.END, spacing=2),
                             ])))

    return ft.Column([
        ft.Container(height=10),
        ft.Container(
            content=ft.Column([
                txt("⚡ Universal Converter", size=22, bold=True, color=ACCENT),
                txt("v2.1 Final · Flet Edition", size=12, color=TEXT2),
            ], horizontal_alignment=ft.CrossAxisAlignment.CENTER),
            alignment=ft.alignment.center, padding=16),

        card(ft.Column([
            txt("📋  Recent Conversions", size=13, bold=True),
            ft.Container(height=8),
            *hist_controls,
        ])),

        card(ft.Column([
            txt("Dependency Status", size=13, bold=True),
            ft.Container(height=8),
            *dep_rows,
        ])),

        card(ft.Column([
            txt("🌐 Flet Edition", size=11, color=TEXT2),
            txt("flet build web --module-name Web", size=10, color=YELLOW, mono=True),
            txt("flet build apk  --module-name Web", size=10, color=YELLOW, mono=True),
        ]), color=CARD2),
    ], scroll=ft.ScrollMode.AUTO, expand=True, spacing=0)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main(page: ft.Page):
    page.title       = "Universal Converter"
    page.theme_mode  = ft.ThemeMode.DARK
    page.bgcolor     = BG
    page.padding     = 0
    page.window.width  = 1100
    page.window.height = 720

    NAV_ITEMS = [
        ("📏", "Units",    lambda: build_units_page(page)),
        ("💱", "Currency", lambda: build_currency_page(page)),
        ("🎵", "Audio",    lambda: build_audio_page(page)),
        ("🎬", "Video",    lambda: build_video_page(page)),
        ("🖼", "Images",   lambda: build_image_page(page)),
        ("📄", "PDF",      lambda: build_pdf_page(page)),
        ("ℹ",  "About",   lambda: build_about_page(page)),
    ]

    page_title = ft.Text("Units Converter", size=15,
                          weight=ft.FontWeight.BOLD, color=TEXT)
    content    = ft.Container(expand=True, padding=ft.padding.all(22),
                               bgcolor=BG, content=ft.Container())
    nav_btns   = []
    nav_col    = ft.Column(spacing=2)

    def show(name, builder, btn):
        for b in nav_btns:
            b.style = ft.ButtonStyle(color=TEXT2, bgcolor=ft.Colors.TRANSPARENT)
        btn.style = ft.ButtonStyle(color=ACCENT, bgcolor=CARD2,
                                    shape=ft.RoundedRectangleBorder(radius=8))
        page_title.value = name + (" Converter" if name != "About" else "")
        content.content  = builder()
        page.update()

    for icon, name, builder in NAV_ITEMS:
        btn = ft.TextButton(
            text=f"  {icon}  {name}",
            style=ft.ButtonStyle(color=TEXT2, bgcolor=ft.Colors.TRANSPARENT,
                                  shape=ft.RoundedRectangleBorder(radius=8)))
        def _cb(e, n=name, b=builder, bt=btn): show(n, b, bt)
        btn.on_click = _cb
        nav_btns.append(btn)
        nav_col.controls.append(
            ft.Container(content=btn, margin=ft.margin.symmetric(horizontal=4)))

    def ind(lbl, ok):
        return ft.Text(f"{'●' if ok else '○'} {lbl}", size=10,
                       color=GREEN if ok else TEXT3)

    sidebar = ft.Container(
        width=200, bgcolor=SURFACE,
        content=ft.Column([
            ft.Container(
                content=ft.Column([
                    txt("⚡ Universal", size=15, bold=True, color=ACCENT),
                    txt("C O N V E R T E R", size=8, color=TEXT3),
                ], horizontal_alignment=ft.CrossAxisAlignment.CENTER),
                bgcolor=CARD, padding=18),
            ft.Container(
                content=ft.Text("CONVERTERS", size=9, color=TEXT3),
                padding=ft.padding.only(left=16, top=14, bottom=4)),
            nav_col,
            ft.Container(expand=True),
            ft.Container(
                content=ft.Column([
                    ind("FFmpeg",  config.HAS_FFMPEG),
                    ind("Pillow",  config.HAS_PIL),
                    ind("MS Office", config.HAS_OFFICE),
                ], spacing=4),
                padding=ft.padding.all(12),
                border=ft.border.only(top=ft.BorderSide(1, BORDER))),
        ], spacing=0, expand=True))

    header = ft.Container(
        bgcolor=SURFACE, height=52,
        content=ft.Row([
            page_title,
            ft.Container(expand=True),
        ], vertical_alignment=ft.CrossAxisAlignment.CENTER),
        padding=ft.padding.symmetric(horizontal=20))

    page.add(ft.Row([
        sidebar,
        ft.Column([header, content], expand=True, spacing=0),
    ], expand=True, spacing=0))

    show("Units", lambda: build_units_page(page), nav_btns[0])
    nav_btns[0].style = ft.ButtonStyle(color=ACCENT, bgcolor=CARD2,
                                        shape=ft.RoundedRectangleBorder(radius=8))
    page.update()


if __name__ == "__main__":
    ft.app(target=main)